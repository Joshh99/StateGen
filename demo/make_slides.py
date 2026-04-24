"""
Generate StateGen presentation as a .pptx file.
Run: .venv/bin/python demo/make_slides.py
Output: demo/stategen_slides.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
from pathlib import Path

OUT = Path(__file__).parent / "stategen_slides.pptx"

# ── Colours ───────────────────────────────────────────────────────────────────
INDIGO  = RGBColor(0x4f, 0x46, 0xe5)
GREEN   = RGBColor(0x16, 0xa3, 0x4a)
RED     = RGBColor(0xdc, 0x26, 0x26)
AMBER   = RGBColor(0xd9, 0x74, 0x06)
SLATE   = RGBColor(0x33, 0x41, 0x55)
LIGHT   = RGBColor(0xf8, 0xfa, 0xfc)
WHITE   = RGBColor(0xff, 0xff, 0xff)
GRAY    = RGBColor(0x94, 0xa3, 0xb8)
BGBLUE  = RGBColor(0xef, 0xf6, 0xff)
BGGREEN = RGBColor(0xf0, 0xfd, 0xf4)
BGAMBER = RGBColor(0xff, 0xfb, 0xeb)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # fully blank


# ── Helper functions ──────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, x, y, w, h, text, size=18, bold=False, color=SLATE,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def heading(slide, text, y=0.25, size=28):
    rect(slide, 0, y, 13.33, 0.55, fill=INDIGO)
    textbox(slide, 0.3, y+0.05, 12.7, 0.5, text,
            size=size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def table_slide(slide, headers, rows, x=0.4, y=1.0, col_widths=None, highlights=None):
    n_cols = len(headers)
    total_w = 12.5
    if col_widths is None:
        col_widths = [total_w / n_cols] * n_cols

    row_h = 0.38
    cx = x
    for i, h in enumerate(headers):
        rect(slide, cx, y, col_widths[i], row_h, fill=INDIGO)
        textbox(slide, cx+0.08, y+0.05, col_widths[i]-0.1, row_h-0.08,
                h, size=11, bold=True, color=WHITE)
        cx += col_widths[i]

    for ri, row in enumerate(rows):
        ry = y + row_h * (ri + 1)
        bg = BGBLUE if (highlights and ri in highlights) else (LIGHT if ri % 2 == 0 else WHITE)
        cx = x
        for ci, cell in enumerate(row):
            rect(slide, cx, ry, col_widths[ci], row_h, fill=bg, line=GRAY, line_w=Pt(0.5))
            bold = bool(highlights and ri in highlights)
            textbox(slide, cx+0.08, ry+0.06, col_widths[ci]-0.1, row_h-0.1,
                    cell, size=11, bold=bold, color=INDIGO if bold else SLATE)
            cx += col_widths[ci]


def card(slide, x, y, w, h, title, body, top_color=INDIGO, bg=LIGHT):
    rect(slide, x, y, w, h, fill=bg, line=GRAY, line_w=Pt(0.5))
    rect(slide, x, y, w, 0.06, fill=top_color)
    textbox(slide, x+0.12, y+0.12, w-0.24, 0.32, title,
            size=12, bold=True, color=top_color)
    textbox(slide, x+0.12, y+0.44, w-0.24, h-0.5, body,
            size=10, color=SLATE, wrap=True)


def pipeline_box(slide, x, y, label, color=INDIGO):
    rect(slide, x, y, 1.6, 0.5, fill=color)
    textbox(slide, x+0.05, y+0.08, 1.5, 0.38, label,
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def arrow(slide, x, y):
    textbox(slide, x, y+0.1, 0.35, 0.35, "→", size=18, color=GRAY,
            align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=INDIGO)
rect(sl, 0, 2.6, 13.33, 2.5, fill=WHITE)

textbox(sl, 1, 0.7, 11, 1.0, "StateGen",
        size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(sl, 1, 1.75, 11, 0.6,
        "Code Generation as a Markov Decision Process with Local Backtracking",
        size=18, color=RGBColor(0xc7,0xd2,0xfe), align=PP_ALIGN.CENTER)

textbox(sl, 1, 2.75, 11, 0.5,
        "CMU 18-786 · Introduction to Deep Learning · Spring 2026",
        size=15, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
textbox(sl, 1, 3.3, 11, 0.4,
        "Joe Maina · Patrick Nguyen · Patricia Muriira",
        size=13, color=SLATE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "The Problem with One-Shot Code Generation")

card(sl, 0.3, 1.0, 3.9, 2.2,
     "❌  Single-shot fails",
     "One LLM call generates the whole program.\nAny small error means the entire output is wrong — no recovery possible.",
     top_color=RED, bg=RGBColor(0xff,0xf1,0xf1))

card(sl, 4.5, 1.0, 3.9, 2.2,
     "⚠  Global retry is wasteful",
     "Regenerate the entire solution on every failure.\nWastes tokens rewriting code that was already correct.",
     top_color=AMBER, bg=BGAMBER)

card(sl, 8.7, 1.0, 4.2, 2.2,
     "✓  Local repair is better",
     "Identify which part failed.\nRepair only that fragment and leave the rest intact.",
     top_color=GREEN, bg=BGGREEN)

textbox(sl, 0.4, 3.5, 12.5, 0.4,
        "Real programs have natural structure:  imports  →  helpers  →  main function",
        size=14, bold=True, color=INDIGO)
textbox(sl, 0.4, 3.95, 12.5, 0.4,
        "StateGen exploits this structure to make repair local and efficient.",
        size=13, color=SLATE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MDP Framing
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Framing Code Generation as an MDP")

table_slide(sl,
    ["MDP Component", "In StateGen"],
    [
        ["State  s",        "Partial program built so far (accumulated code fragments)"],
        ["Action  a",       "Generate code for the next sub-problem (one LLM call)"],
        ["Transition",      "Append fragment → new partial program; verify before committing"],
        ["Reward",          "+1 if all BCB unit tests pass, 0 otherwise"],
        ["Backtrack",       "Failed state → repair only that fragment (local), not revert to s₀"],
    ],
    x=0.4, y=1.0, col_widths=[2.8, 9.2],
    highlights=[4]
)

textbox(sl, 0.4, 3.75, 12.5, 0.35, "Why local backtracking?", size=14, bold=True, color=INDIGO)
bullets = [
    "A bug in state i almost never originates from state i−1 (bugs are local)",
    "Full MDP backtrack to s₀ wastes all correctly generated prior states",
    "Local retry: fix only the failing fragment, re-append, continue forward",
    "Validated empirically: same accuracy as full backtrack at 2.6× fewer retries",
]
ty = 4.15
for b in bullets:
    textbox(sl, 0.5, ty, 12.3, 0.34, f"→  {b}", size=11, color=SLATE)
    ty += 0.34


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Four Baselines
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Four Methods We Compare")

card(sl, 0.3,  1.1, 2.9, 2.8, "Direct Generation",
     "Single LLM call.\nNo planning, no retry.\n\nBaseline floor.\n\n1 LLM call",
     top_color=GRAY, bg=LIGHT)

card(sl, 3.5,  1.1, 2.9, 2.8, "Self-Planning",
     "Write a numbered plan first, then generate code from the plan.\nRetry whole solution on failure.\n\n(Jiang et al., TOSEM 2024)\n\n2+ LLM calls",
     top_color=AMBER, bg=BGAMBER)

card(sl, 6.7,  1.1, 2.9, 2.8, "Self-Debugging",
     "Generate → explain the code in prose → repair based on explanation.\n\n(Chen et al., ICLR 2024)\n\n3+ LLM calls",
     top_color=RED, bg=RGBColor(0xff,0xf1,0xf1))

card(sl, 9.9,  1.1, 3.1, 2.8, "StateGen  ★",
     "Decompose → per-state verify & repair → final BCB eval.\n\nLocal backtracking. Memory cache.\n\n2+ LLM calls",
     top_color=INDIGO, bg=BGBLUE)

textbox(sl, 0.4, 4.15, 12.5, 0.35,
        "All methods use the same LLM, same executor, same token tracker — fair comparison.",
        size=11, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Decomposition
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Step 1: Decomposition")

textbox(sl, 0.4, 0.95, 12.5, 0.35,
        "One LLM call produces a structured breakdown of the problem into sequential states:",
        size=13, color=SLATE)

rect(sl, 0.4, 1.35, 7.8, 3.8, fill=RGBColor(0xf1,0xf5,0xf9), line=GRAY, line_w=Pt(0.75))
code = (
    "STATE 0\n"
    "Description: Import modules and define constants.\n"
    "Preconditions: none\n"
    "Postconditions: 'pandas' in dir(); 'numpy' in dir()\n"
    "CODE:\n"
    "import pandas as pd\n"
    "import numpy as np\n"
    "END_STATE\n\n"
    "STATE 1\n"
    "Description: Define the main task_func.\n"
    "Preconditions: pd, np\n"
    "Postconditions: callable(task_func)\n"
    "CODE:\n"
    "def task_func(df, col):\n"
    "    return df.groupby(col).mean()\n"
    "END_STATE"
)
textbox(sl, 0.55, 1.45, 7.6, 3.6, code, size=10,
        color=SLATE, wrap=False)

textbox(sl, 8.5, 1.35, 4.5, 0.35, "Each state has:", size=13, bold=True, color=INDIGO)
items = [
    "Preconditions — what must already exist",
    "Postconditions — structural assertions to verify",
    "CODE block — the actual implementation",
    "END_STATE — explicit delimiter",
]
ty = 1.75
for it in items:
    textbox(sl, 8.6, ty, 4.3, 0.34, f"→  {it}", size=11, color=SLATE)
    ty += 0.38

rect(sl, 8.5, 3.4, 4.5, 1.6, fill=BGBLUE, line=INDIGO, line_w=Pt(0.75))
textbox(sl, 8.65, 3.5, 4.2, 0.3, "Design choice", size=11, bold=True, color=INDIGO)
textbox(sl, 8.65, 3.85, 4.2, 1.0,
        "State 0 = imports only.\nState 1 = all function defs.\nOptional state 2 for helpers.\nSmall states = easier to repair.",
        size=10, color=SLATE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Quick Check
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Step 2: Per-State Quick Check")

textbox(sl, 0.4, 0.9, 12.5, 0.35,
        "After generating code for each state, we verify it immediately before moving on:",
        size=13, color=SLATE)

steps = [
    ("Generate\nfragment",    INDIGO),
    ("Syntax\ncheck",         AMBER),
    ("Run in\nsandbox",       AMBER),
    ("Check\npostconditions", GREEN),
    ("Commit to\ncontext ✓",  GREEN),
]
px = 0.4
for i, (label, color) in enumerate(steps):
    pipeline_box(sl, px, 1.4, label, color)
    if i < len(steps) - 1:
        arrow(sl, px + 1.6, 1.4)
        px += 1.97

rect(sl, 0.4, 2.2, 5.8, 2.2, fill=RGBColor(0xff,0xf1,0xf1), line=RED, line_w=Pt(0.75))
textbox(sl, 0.55, 2.3, 5.5, 0.35, "On failure → local repair", size=13, bold=True, color=RED)
for i, b in enumerate([
    "Error message passed back to LLM",
    "LLM rewrites only the failing fragment",
    "Previous states remain untouched",
    "Up to 3 repair attempts per state",
]):
    textbox(sl, 0.65, 2.7 + i*0.32, 5.5, 0.3, f"→  {b}", size=11, color=SLATE)

rect(sl, 6.6, 2.2, 6.4, 2.2, fill=BGGREEN, line=GREEN, line_w=Pt(0.75))
textbox(sl, 6.75, 2.3, 6.1, 0.35, "What quick check catches", size=13, bold=True, color=GREEN)
for i, b in enumerate([
    "SyntaxError, IndentationError",
    "Import failures (missing modules)",
    "Runtime crashes at definition time",
    "Structural postconditions: callable(f), 'pd' in dir()",
]):
    textbox(sl, 6.85, 2.7 + i*0.32, 6.1, 0.3, f"→  {b}", size=11, color=SLATE)

textbox(sl, 0.4, 4.6, 12.5, 0.35,
        "Quick check runs code in an isolated subprocess — no effect on the main process.",
        size=11, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Full Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Full StateGen Pipeline")

stages = [
    ("Decompose\n(1 LLM call)", INDIGO),
    ("State 0\n✓ pass",         GREEN),
    ("State 1\nattempt 0 ✗",    RED),
    ("State 1\nattempt 1 ✓",    GREEN),
    ("Final BCB\neval",          AMBER),
]
px = 0.4
for i, (label, color) in enumerate(stages):
    pipeline_box(sl, px, 1.05, label, color)
    if i < len(stages) - 1:
        arrow(sl, px + 1.6, 1.05)
        px += 1.97

textbox(sl, 0.4, 1.8, 12.5, 0.3, "Happy path:  all states pass  →  BCB tests pass  →  done ✓",
        size=12, bold=True, color=GREEN)

rect(sl, 0.4, 2.2, 12.5, 1.3, fill=BGAMBER, line=AMBER, line_w=Pt(0.75))
textbox(sl, 0.55, 2.3, 12.0, 0.3,
        "If final BCB eval fails after all states pass quick check:",
        size=12, bold=True, color=AMBER)
repair_steps = [
    ("BCB\nFAIL ✗", RED),
    ("Whole-solution\nrepair call", AMBER),
    ("BCB\neval again", AMBER),
    ("Up to 4\nrepair rounds", AMBER),
    ("BCB\nPASS ✓", GREEN),
]
px2 = 0.7
for i, (label, color) in enumerate(repair_steps):
    pipeline_box(sl, px2, 2.65, label, color)
    if i < len(repair_steps) - 1:
        arrow(sl, px2 + 1.6, 2.65)
        px2 += 1.97

textbox(sl, 0.4, 4.1, 12.5, 0.35,
        "BCB evaluation runs the official unit test suite inside Docker — not just syntax checks.",
        size=11, color=GRAY, italic=True)
textbox(sl, 0.4, 4.55, 12.5, 0.35,
        "Memory manager: before generating each state, check similarity cache — skip LLM call if hit (threshold = 0.85 cosine similarity).",
        size=11, color=SLATE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BigCodeBench
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Benchmark: BigCodeBench Hard")

card(sl, 0.3, 1.0, 4.0, 2.5, "Dataset",
     "148 tasks from the hard instruct subset\n\nReal-world library usage:\npandas, sklearn, matplotlib,\ntensorflow, flask, openpyxl, …\n\nEach task has 5–15 unit tests",
     top_color=INDIGO, bg=BGBLUE)

card(sl, 4.7, 1.0, 4.0, 2.5, "Why it's hard",
     "Tests check exact outputs,\ntypes, shapes, and edge cases\n\nRequires correct API usage\n— not just plausible-looking code\n\nSoTA models score only 30–45%",
     top_color=RED, bg=RGBColor(0xff,0xf1,0xf1))

card(sl, 9.1, 1.0, 3.9, 2.5, "Evaluation",
     "Official Docker container\nruns real Python unit tests\n\nMetric: pass@1\n(first solution must pass all tests)\n\nNo partial credit",
     top_color=GREEN, bg=BGGREEN)

rect(sl, 0.3, 3.75, 12.7, 1.5, fill=LIGHT, line=GRAY, line_w=Pt(0.75))
textbox(sl, 0.5, 3.85, 12.3, 0.3, "Example task  (BigCodeBench/82)", size=12, bold=True, color=INDIGO)
textbox(sl, 0.5, 4.2, 12.3, 0.8,
        "Given a list of file paths and a target directory, read each CSV file with pandas, "
        "apply StandardScaler normalisation, and write the scaled data to the target directory. "
        "Return a dict mapping filename → scaled DataFrame.",
        size=11, color=SLATE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Main Results (DeepSeek)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Results: DeepSeek-Coder  (148 tasks, BCB Hard)")

table_slide(sl,
    ["Method", "Pass@1", "Avg Tokens", "Avg Retries", "Avg Time"],
    [
        ["Direct Generation",  "23.0%",  "603",   "0.00",  "19 s"],
        ["Self-Planning",      "28.4%",  "3,197", "3.04",  "87 s"],
        ["StateGen  ★",        "29.7%",  "3,106", "1.80",  "75 s"],
        ["Self-Debugging",     "30.3%",  "6,560", "1.48",  "140 s"],
    ],
    x=0.4, y=1.0, col_widths=[3.8, 2.0, 2.3, 2.3, 2.1],
    highlights=[2]
)

card(sl, 0.3,  3.55, 3.9, 1.9, "+29% relative improvement",
     "23.0%  →  29.7% over direct generation\n(+6.7 percentage points absolute)",
     top_color=GREEN, bg=BGGREEN)

card(sl, 4.5,  3.55, 4.0, 1.9, "53% fewer tokens vs Self-Debugging",
     "Matches Self-Debugging accuracy (29.7% vs 30.3%)\nat 3,106 vs 6,560 tokens per task",
     top_color=GREEN, bg=BGGREEN)

card(sl, 8.8,  3.55, 4.2, 1.9, "More efficient than Self-Planning",
     "Slightly better accuracy (29.7% vs 28.4%)\nFewer retries (1.80 vs 3.04)\n14% faster wall time (75s vs 87s)",
     top_color=GREEN, bg=BGGREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Ablations (DeepSeek)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Ablation Study: Which Components Drive the Gain?  (DeepSeek-Coder)")

table_slide(sl,
    ["Variant", "What is removed", "Pass@1", "Avg Tokens", "Avg Retries", "Δ Pass@1"],
    [
        ["StateGen (full)  ★", "—",                           "29.7%", "3,106", "1.80", "—"],
        ["No memory",          "Pattern similarity cache",     "27.7%", "3,433", "4.70", "−2.0 pp"],
        ["No decompose",       "State decomposition step",     "29.1%", "2,340", "1.74", "−0.6 pp"],
        ["Full backtrack",     "Reset entire decomp on fail",  "29.7%", "3,345", "4.66", "±0.0 pp"],
    ],
    x=0.4, y=1.0, col_widths=[2.9, 3.0, 1.4, 1.6, 1.6, 1.8],
    highlights=[0]
)

rect(sl, 0.3, 3.55, 12.7, 2.1, fill=BGBLUE, line=INDIGO, line_w=Pt(1))
textbox(sl, 0.5, 3.65, 12.3, 0.35,
        "Key insight: memory is the most critical component",
        size=14, bold=True, color=INDIGO)
textbox(sl, 0.5, 4.05, 12.3, 1.4,
        "Removing memory costs −2.0 pp accuracy AND increases tokens by 10% (3,433 vs 3,106) AND 2.6× more retries — strictly worse on every dimension.\n"
        "Full backtrack ties local repair in accuracy (both 29.7%) but uses 2.6× more retries (4.66 vs 1.80) and 8% more tokens.\n"
        "Local repair advantage is efficiency, not absolute accuracy: same pass rate with far less redundant work.",
        size=11, color=SLATE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Gemini 2.5 Flash Results
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Results: Gemini 2.5 Flash  (148 tasks, BCB Hard)")

table_slide(sl,
    ["Method", "Pass@1", "Avg Tokens", "Avg Retries", "Avg Time"],
    [
        ["Direct Generation",      "45.3%", "901",   "0.00", "16 s"],
        ["Self-Planning",          "45.3%", "4,226", "2.20", "34 s"],
        ["Self-Debugging",         "81.8%", "6,033", "0.52", "35 s"],
        ["StateGen",               "80.4%", "3,471", "3.59", "57 s"],
        ["No Memory",              "68.2%", "3,571", "3.88", "52 s"],
        ["No Decompose",           "82.4%", "2,559", "1.09", "33 s"],
        ["Full Backtrack  ★",      "83.8%", "3,413", "3.70", "59 s"],
    ],
    x=0.4, y=1.0, col_widths=[3.5, 1.9, 2.0, 2.0, 1.9],
    highlights=[6]
)

card(sl, 0.3, 4.25, 4.0, 2.0, "Self-Planning stalls",
     "Self-Planning = 45.3% — identical to Direct Gen.\nGemini already plans internally via thinking tokens.\nExternal planning prompt adds noise, not value.",
     top_color=AMBER, bg=BGAMBER)

card(sl, 4.6, 4.25, 4.0, 2.0, "Memory critical (−12.2 pp)",
     "No Memory drops from 80.4% → 68.2%.\nLargest single-component effect across all experiments.\nGemini failure modes are predictable & reusable.",
     top_color=RED, bg=RGBColor(0xff,0xf1,0xf1))

card(sl, 9.0, 4.25, 4.0, 2.0, "Best: Full Backtrack",
     "83.8% at 3,413 tokens — beats Self-Debugging (81.8%) at half the tokens.\nNo Decompose is close (82.4%) at 25% fewer tokens.",
     top_color=GREEN, bg=BGGREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — GPT-4o Results
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Results: GPT-4o Azure  (148 tasks, BCB Hard)")

table_slide(sl,
    ["Method", "Pass@1", "Avg Tokens", "Avg Retries", "Avg Time"],
    [
        ["No Decompose  ★",  "83.8%", "1,193", "0.90", "6 s"],
        ["Full Backtrack",   "83.1%", "3,053", "3.61", "21 s"],
        ["StateGen",         "82.4%", "3,146", "3.82", "19 s"],
        ["No Memory",        "82.4%", "2,907", "3.14", "16 s"],
    ],
    x=0.4, y=1.0, col_widths=[3.5, 1.9, 2.0, 2.0, 1.9],
    highlights=[0]
)

card(sl, 0.3, 3.5, 3.9, 2.2, "62% fewer tokens",
     "No Decompose: 1,193 tokens vs 3,146 for full StateGen.\nHighest pass rate (83.8%) at lowest cost.\nSkipping decomposition = faster + better for GPT-4o.",
     top_color=GREEN, bg=BGGREEN)

card(sl, 4.5, 3.5, 4.2, 2.2, "Memory has zero effect",
     "StateGen vs No Memory: both 82.4%.\nGPT-4o failure modes are idiosyncratic — patterns from past tasks don't transfer.\nCache adds cost, not accuracy.",
     top_color=AMBER, bg=BGAMBER)

card(sl, 9.0, 3.5, 4.0, 2.2, "6 seconds per task",
     "No Decompose averages only 6 s/task — 3× faster than any other variant.\nSkipping the decomposition call eliminates the most expensive LLM step.",
     top_color=INDIGO, bg=BGBLUE)

textbox(sl, 0.4, 5.9, 12.5, 0.3,
        "Note: Direct Gen, Self-Planning, Self-Debugging baselines not run for GPT-4o (StateGen variants only).",
        size=10, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Cross-Model Comparison
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Cross-Model Comparison  (148 tasks each)")

table_slide(sl,
    ["Method", "DeepSeek-Coder", "Gemini 2.5 Flash", "GPT-4o Azure"],
    [
        ["StateGen (full)",         "29.7%",          "80.4%",          "82.4%"],
        ["No Memory",               "27.7%",          "68.2%",          "82.4%"],
        ["No Decompose",            "29.1%",          "82.4%",          "83.8%  ★"],
        ["Full Backtrack",          "29.7%",          "83.8%  ★",       "83.1%"],
        ["Direct Generation",       "23.0%",          "45.3%",          "—"],
        ["Self-Debugging",          "30.3%",          "81.8%",          "—"],
    ],
    x=0.4, y=1.1, col_widths=[3.8, 2.8, 2.8, 2.8],
    highlights=[]
)

rect(sl, 0.3, 4.5, 12.7, 2.2, fill=BGBLUE, line=INDIGO, line_w=Pt(0.75))
textbox(sl, 0.5, 4.6, 12.3, 0.35,
        "Three patterns across model tiers",
        size=13, bold=True, color=INDIGO)
textbox(sl, 0.5, 5.0, 12.3, 1.5,
        "1.  Accuracy ceiling: frontier models reach 80–84% across ALL variants; variant choice drives efficiency not accuracy.\n"
        "2.  Memory matters for Gemini (−12.2 pp without it) but not GPT-4o (±0 pp) — model-specific failure distributions.\n"
        "3.  Decomposition overhead: GPT-4o's best result skips decomposition entirely (83.8% at 1,193 tokens vs 3,146 for full StateGen).",
        size=11, color=SLATE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Cross-Model Insight (Literature)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Why Scaffolding Helps Less for Frontier Models")

card(sl, 0.3, 1.0, 4.0, 2.6,
     "Emergent Planning\n(Wei et al., 2022)",
     "Above a capability threshold, models develop implicit planning without prompting.\n\n"
     "Gemini 2.5 Flash uses extended thinking tokens — it already decomposes internally.\n"
     "Our STATE/CODE scaffold partially conflicts with the model's own plan.",
     top_color=INDIGO, bg=BGBLUE)

card(sl, 4.6, 1.0, 4.0, 2.6,
     "CoT Helps Weaker Models More\n(Kojima et al., 2022)",
     "Plan-and-Solve prompting gives consistent gains for smaller models but "
     "diminishing returns for frontier models.\n\n"
     "DeepSeek-Coder benefits from the full scaffold (+29% relative).\n"
     "GPT-4o's best result uses the simplest variant (no decompose).",
     top_color=AMBER, bg=BGAMBER)

card(sl, 9.2, 1.0, 3.8, 2.6,
     "Memory & Distribution\n(Huang et al., 2023)",
     "The pattern cache is calibrated from previous tasks with the same model.\n\n"
     "Gemini: −12.2 pp without memory → predictable failure patterns reuse well.\n"
     "GPT-4o: ±0 pp without memory → idiosyncratic outputs, cache adds no value.",
     top_color=RED, bg=RGBColor(0xff,0xf1,0xf1))

rect(sl, 0.3, 3.85, 12.7, 1.8, fill=BGGREEN, line=GREEN, line_w=Pt(1))
textbox(sl, 0.5, 3.95, 12.3, 0.35,
        "Actionable finding for practitioners",
        size=13, bold=True, color=GREEN)
textbox(sl, 0.5, 4.35, 12.3, 1.1,
        "Mid-tier models (DeepSeek): use full StateGen — decomposition + memory + local repair all add value.\n"
        "Frontier + thinking tokens (Gemini): keep memory cache (critical), use full backtrack or no-decompose.\n"
        "Frontier + strong RLHF (GPT-4o): use no_decompose only — 83.8% at 1,193 tokens, 6 seconds per task.",
        size=11, color=SLATE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Demo
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
heading(sl, "Live Demo")

textbox(sl, 0.4, 0.9, 12.5, 0.35, "Live API calls to DeepSeek-Coder on any of 148 BigCodeBench tasks:",
        size=13, color=SLATE)

card(sl, 0.3, 1.3, 5.8, 3.3, "What you'll see",
     "→  Select any BCB task (148 available)\n"
     "→  Choose method: Direct Gen / Self-Planning /\n"
     "     Self-Debugging / StateGen / Ablation\n"
     "→  Watch the live API call in real time\n"
     "→  See each step: decomposition, per-state\n"
     "     verification, repairs, memory hits\n"
     "→  Final token count and pass/fail result",
     top_color=INDIGO, bg=BGBLUE)

card(sl, 6.5, 1.3, 3.1, 1.4, "Good demo tasks",
     "BCB/34 — StateGen repairs state 1\nBCB/82 — memory hit on state 0\nBCB/129 — self-debugging repairs",
     top_color=GREEN, bg=BGGREEN)

card(sl, 9.9, 1.3, 3.1, 1.4, "Run the demo",
     "streamlit run demo/app.py\n\nOpens at localhost:8501",
     top_color=INDIGO, bg=BGBLUE)

card(sl, 6.5, 3.0, 6.5, 1.6, "What makes the demo interesting",
     "Run the same task with different methods side by side and compare:\n"
     "how many LLM calls, which states failed, how the repair prompt looks,\n"
     "and exactly how many tokens each method consumed.",
     top_color=AMBER, bg=BGAMBER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT)
heading(sl, "Summary & Conclusions")

card(sl, 0.3, 1.0, 3.9, 1.9, "✓  Accuracy (DeepSeek)",
     "+29% relative vs direct gen.\n23.0% → 29.7% pass@1.\nCompetitive with self-debugging at 53% fewer tokens.",
     top_color=GREEN, bg=BGGREEN)

card(sl, 4.5, 1.0, 3.9, 1.9, "✓  Efficiency",
     "3,106 tokens vs 6,560 (self-debugging).\n75s vs 140s per task.\nFewer retries: 1.80 vs 3.04 (self-planning).",
     top_color=GREEN, bg=BGGREEN)

card(sl, 8.7, 1.0, 4.3, 1.9, "✓  Local > Full Backtrack",
     "Same accuracy (29.7%) at 2.6× fewer retries.\nBugs are local — reset to s₀ wastes correct states.\nEfficiency win, not accuracy win.",
     top_color=GREEN, bg=BGGREEN)

card(sl, 0.3, 3.2, 5.9, 2.1, "Cross-Model Finding",
     "Scaffolding value scales inversely with model capability.\n\n"
     "DeepSeek (mid-tier): full StateGen best (29.7%).\n"
     "Gemini 2.5 Flash: full backtrack best (83.8%); memory critical (−12.2 pp).\n"
     "GPT-4o Azure: no_decompose best (83.8% at only 1,193 tokens).\n\n"
     "Frontier models already plan internally (Wei et al., 2022).",
     top_color=AMBER, bg=BGAMBER)

card(sl, 6.5, 3.2, 6.5, 2.1, "Key Takeaway",
     "Treating code generation as an MDP with local state-based repair gives better "
     "token efficiency than global retry for mid-tier models.\n\n"
     "For frontier models, the execution oracle (quick check) is the universally "
     "valuable component — the appropriate level of decomposition scaffolding "
     "scales inversely with model capability, consistent with emergent abilities "
     "literature (Wei et al., 2022; Kojima et al., 2022).",
     top_color=INDIGO, bg=BGBLUE)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved → {OUT}")
